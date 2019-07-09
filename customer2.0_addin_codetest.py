#' Business Development
#' Confidential & Proprietary
#' Not for Distribution (except to authorized persons)
#' DRAFT
#' Customer reporting metrics presentation additions
#' Purpose: This script creates a Data Metrics presentation for hospice
#'

# import modules
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.chart import XL_LEGEND_POSITION
from pptx.enum.chart import XL_TICK_MARK
from pptx.enum.chart import XL_LABEL_POSITION
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Inches
from pptx.util import Pt
import pyodbc
import cx_Oracle
import csv
import os
import pandas as pd
import numpy as np
import json
import datetime as dt

# create slide templates

templateX = 'CHHA Rehospitalization Trend'


summaryX= 'CHHA Rehospitalization Trends'


''''''''''''''''''''''''''''''''''''''''''''''''''''''''''''''''''''''''''''''
# chart size and location variables:
chart_height= Inches(4.99)
chart_width= Inches(8.85)
chart_x= Inches(.71)
chart_y= Inches(1.25)

# n totals size & location variables:
n_height= Inches(.3)
n_width= Inches(7.9)
n_x= Inches(1.41)
n_y= Inches(1.6)

# comments box size & location variables:
c_left= Inches(1.37)
c_top= Inches(6.17)
c_width= Inches(7.52)
c_height= Inches(.34)

# other text size & location variables:
oleft= Inches(.19)
otop=Inches(6.79)
owidth=Inches(7.72)
oheight=Inches(.77)


'''functions'''''''''''''''''''''''''''''''''''''''''''''''''''''''''''''''''''''
# This section uses functions to build the different pieces of the presentation

# create the comments box for all slides


def comments(x):
    txBox = x.shapes.add_textbox(c_left, c_top, c_width, c_height)
    tf = txBox.text_frame
    tf.word_wrap= True
    p = tf.paragraphs[0]
    p.alignment=PP_ALIGN.CENTER
    run = p.add_run()
    run.text='Comments'
    font = run.font
    font.size = Pt(14)
    fill=txBox.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(247,247,247)
    txBox.line.fill.solid()
    txBox.line.dash_style=MSO_LINE.SOLID
    return


# creating line chart.


def chart_layout_line_p(x,chart_data,y):
    graphic_frame=x.shapes.add_chart(
        XL_CHART_TYPE.LINE, chart_x, chart_y, chart_width, chart_height, chart_data)
    chart=graphic_frame.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size= Pt(12)
    chart.chart_title.text_frame.text=y
    category_axis= chart.category_axis
    category_axis.tick_labels.font.size = Pt(12)
    value_axis = chart.value_axis
    value_axis.minimum_scale = 0
    value_axis.has_major_gridlines = False
    tick_labels = value_axis.tick_labels
    tick_labels.font.size = Pt(12)
    plot = chart.plots[0]
    plot.has_data_labels = True
    data_labels = plot.data_labels
    data_labels.font.size = Pt(12)
    return


#creating clustered column chart

def chart_layout_column(x,chart_data,y):
    graphic_frame=x.shapes.add_chart(
       XL_CHART_TYPE.COLUMN_CLUSTERED, chart_x, chart_y, chart_width, chart_height, chart_data)
    chart=graphic_frame.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size= Pt(12)
    chart.chart_title.text_frame.text=y
    category_axis= chart.category_axis
    category_axis.tick_labels.font.size = Pt(12)
    value_axis = chart.value_axis
    value_axis.has_major_gridlines = False
    tick_labels = value_axis.tick_labels
    tick_labels.font.size = Pt(12)
    plot = chart.plots[0]
    plot.has_data_labels = True
    data_labels = plot.data_labels
    data_labels.font.size = Pt(12)
    return

#Crearting pie-chart

def chart_layout_pie(x,chart_data,y):
    chart = x.shapes.add_chart(
    XL_CHART_TYPE.PIE, chart_x, chart_y, chart_width, chart_height, chart_data
    ).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.plots[0].has_data_labels = True
    data_labels = chart.plots[0].data_labels
    data_labels.number_format = '0.0%'
    data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
    data_labels.font.size = Pt(12)





def formatter(x):
    df = pd.DataFrame(x)
    count = int(df.count())
    if count > 1:
        return str(tuple(x))
    else:
        return '('+str(x)[1:-1]+')'

"""
Database Connections
"""

local_settings='C:\\Users\\'+os.getlogin()+'\\Documents\\Python Scripts\\'
with open(local_settings+'blank_settings.json') as json_data:
    json_pw=json.load(json_data)

def connections(x):
    if x=="hchb":
        return pyodbc.connect(driver="{SQL Server}", server='HCHBDBW0301', database='HCHB',
                        uid=json_pw['password'][0]['user'],
                       pwd=json_pw['password'][0]['password'])
    else:
        return cx_Oracle.connect(json_pw['password'][1]['user'],
                         json_pw['password'][1]['password'],
                         'RCPROD1')
    return


''''''''''''''''''''''''''''''''''''''''''''''''''''''''''''''''''''''''''''''''

"""
Customer Referral Names
"""

referral_nm = 'MT SINAI BETH ISRAEL', 'MT SINAI BROOKLYN', 'MT SINAI HOSP OF QUEENS', 'MT SINAI HOSPITAL', 'MT SINAI WEST', 'MT SINAI ST LUKES'


q="""
SELECT *
FROM PETER_D.CUSTOMER_REPORTING2
WHERE PARENT in"""+str(referral_nm)

q_cancel="""
SELECT *
FROM PETER_D.XWALK_CANCEL_CODE
"""

q_payor="""
SELECT *
FROM PETER_D.XWALK_PAYOR_TYPE
"""
con=connections("pcrs")

referral_xwalk=pd.read_sql(q, con=con)
cancel_xwalk=pd.read_sql(q_cancel, con=con)
payor_xwalk=pd.read_sql(q_payor, con=con)

con.close()

referral_id=formatter(referral_xwalk['ID'])
df_locations=referral_xwalk[['PARENT','CUSTOMER_HOSPITAL','SYSTEM']].drop_duplicates()
system=df_locations['SYSTEM'].drop_duplicates().to_string(index=False)

input_start='2018-01-01'
input_stop='2019-04-01'

start=pd.to_datetime(input_start)
stop=pd.to_datetime(input_stop)

o_start=start.strftime('%d-%b-%y')
o_stop=stop.strftime('%d-%b-%y')

b_reporting=start.to_period("Q")
e_reporting=(stop-dt.timedelta(days=1)).to_period("Q")

start_30=(start- dt.timedelta(days=30)).strftime('%Y-%m-%d')
stop_30=(stop- dt.timedelta(days=30)).strftime('%Y-%m-%d')

o_start30 = (start-dt.timedelta(days=30))
o_stop30 = (stop-dt.timedelta(days=30))

o_start_30 = o_start30.strftime('%d-%b-%y')
o_stop_30 = o_stop30.strftime('%d-%b-%y')



''''''''''''''''''''''''''''''''''''''''''''''''''''''''''''''''''''''''''''''''
# LOAD DATA
# Admissions
# PCRS

pcrs_q="""
Select c.case_num, C.admission_date, c.referral_ID, pp.SAS_RPT_CAT1_LBL, 'HOME HEALTH' as LOB, R.REG_LABEL as Region
from MS_OWNER.VW_CASE_FACTS c
left join MS_OWNER.DIM_REGION R
on C.BOROUGH = R.REG_CODE
left join MS_OWNER.DIM_PAYOR_PLAN pp
on PP.CD=C.PAYOR_PRIMARY
where c.admission_date >= """+"'"+o_start+"'""""
and c.admission_date < """+"'"+o_stop+"'""""
and c.referral_ID in"""+str(referral_id)+"""
and C.COMPANY_CODE='VNS'
and C.REFERRAL_TYPE='A'
"""
con=connections("pcrs")
df_pcrs=pd.read_sql(pcrs_q, con=con)
con.close()


# HCHB
hchb_q="""
Select c.epi_id, c.epi_socDate, c.epi_referralfaid, pt.pt_desc, sl.sl_desc as LOB, bg.bg_description as Region
from HCHB.dbo.client_episodes_all c
left join CLIENT_EPISODE_FS fs
on c.epi_id=fs.cefs_epiid
left join PAYOR_SOURCES p
on p.ps_id=fs.cefs_psid
left join PAYOR_TYPES pt
on pt.pt_id=fs.cefs_ptid
left join HCHB.dbo.service_lines sl
on c.epi_slid=sl.sl_id
left join [VNSNY_BI].[dbo].[VW_CHHA_BRANCH_GROUPS_TEAMS] b
on b.branch_code=c.epi_branchcode and b.team_id=c.epi_teamid
left join [VNSNY_BI].[dbo].[VW_BRANCH_GROUPS] bg
on bg.bg_id=b.BGB_BGID
where c.epi_AdmitType not in('RECERTIFICATION', 'BEREAVEMENT')
and c.epi_status not in('deleted','PENDING','NON-ADMIT')
and c.epi_socDate >= """+"'"+input_start+"'""""
and c.epi_socDate < """+"'"+input_stop+"'""""
and fs.cefs_ps='P'
and fs.cefs_active='Y'
and c.epi_branchcode != 'COR'
and c.epi_referralfaid in"""+formatter(referral_xwalk.loc[referral_xwalk['SRC'] == 'HCHB', 'ID'])+"""
"""

con=connections("hchb")
df_hchb=pd.read_sql(hchb_q, con=con)
con.close()

# Combined data
df_pcrs['src']='pcrs'
df_hchb['src']='hchb'
df_pcrs.columns=('case_id', 'adm_date', 'referral_id', 'payor_cat', 'lob', 'region','src')
df_hchb.columns=('case_id', 'adm_date', 'referral_id', 'payor_cat', 'lob', 'region','src')
df_hchb['referral_id']=df_hchb['referral_id'].astype(str)
frames=[df_pcrs, df_hchb]
df_adm=pd.concat(frames)
df_adm=df_adm.merge(referral_xwalk, left_on='referral_id', right_on='ID', how='left')
df_adm=df_adm.merge(payor_xwalk, left_on='payor_cat',right_on='PAYOR_TYPE', how='left')
df_adm['date']=df_adm['adm_date'].dt.strftime('%Y-%m-01')
df_adm['quarter']=pd.to_datetime(df_adm['date']).dt.to_period("Q")
df_adm['region'] = df_adm['region'].str.strip()
df_adm = df_adm.replace({'region': {'CHHA MANHATTAN': 'Manhattan', 'CHHA STATEN ISLAND' : 'Staten Island', 'CHHA WESTCHESTER' : 'Westchester', 'CHHA BROOKLYN' : 'Brooklyn', 'CHHA BRONX' : 'Bronx', 'CHHA NASSAU SUFFOLK' : 'Suffolk', 'CHHA QUEENS' : 'Queens'}})
df_adm = df_adm.replace({'region': {'HOSPICE MANHATTAN' : 'Manhattan', 'HOSPICE BROOKLYN' : 'Brooklyn', 'HOSPICE BRONX' : 'Bronx', 'HOSPICE QUEENS' : 'Queens', 'HOS STATEN ISLAND N' : 'Staten Island'}})



#hchb rehosp (numerator)

hchb_q1=(
"""
select sum(e.HOSPITALIZATION_30_DAY)  as numerator,CAST(year(epa.epi_SocDate+30) AS char(4)) + '-Q' +
           CAST(CEILING(CAST(month(epa.epi_SocDate+30) AS decimal(9,2)) / 3) AS char(1)) as Quarter, epa.epi_ReferralFaId
from  HCHB.dbo.CLIENT_EPISODES_ALL epa
         left outer join   VNSNY_BI.dbo.VW_EPISODE_FLAGS e
            on    (e.epi_id = epa.epi_id)
where   (epa.epi_slid in (1)
and epa.epi_status not in ('DELETED', 'NON-ADMIT', 'PENDING')
and epa.epi_admitType not in ('RECERTIFICATION')
and epa.epi_ReferralFaId in """+str(tuple(referral_xwalk.loc[referral_xwalk['SRC']=='HCHB','ID']))+"""
and CAST(epa.epi_SocDate+30 as date)>= """+"'"+input_start+"'""""
and CAST(epa.epi_SocDate+30 as date)< """+"'"+input_stop+"'"""")
group by CAST(year(epa.epi_SocDate+30) AS char(4)) + '-Q' +
CAST(CEILING(CAST(month(epa.epi_SocDate+30) AS decimal(9,2)) / 3) AS char(1)), epa.epi_ReferralFaId
""")


con=connections("hchb")
df1=pd.read_sql(hchb_q1,con=con)
con.close()


df1['epi_ReferralFaId'] = df1['epi_ReferralFaId'].astype(str)
df1=df1.merge(referral_xwalk, left_on='epi_ReferralFaId', right_on='ID', how='left')
df1=df1.groupby(['Quarter', 'PARENT'])['numerator'].sum().reset_index()


#hchb admissions(denominator)

hchb_q2=("""
select   count(epa.EPI_PAID)  denominator, CAST(year(epa.epi_SocDate+30) AS char(4)) + '-Q' +
           CAST(CEILING(CAST(month(epa.epi_SocDate+30) AS decimal(9,2)) / 3) AS char(1)) as Quarter, epa.epi_ReferralFaId
from         HCHB.dbo.CLIENT_EPISODES_ALL   epa
where   (epa.epi_slid in (1)
and epa.epi_ReferralFaId in """+str(tuple(referral_xwalk.loc[referral_xwalk['SRC']=='HCHB','ID']))+"""
and epa.epi_status not in ('DELETED', 'NON-ADMIT', 'PENDING')
and epa.epi_admitType not in ('RECERTIFICATION')
and CAST(epa.epi_SocDate+30 as date) >= """+"'"+input_start+"'""""
and CAST(epa.epi_SocDate+30 as date) < """+"'"+input_stop+"'"""")
group by CAST(year(epa.epi_SocDate+30) AS char(4)) + '-Q' +
CAST(CEILING(CAST(month(epa.epi_SocDate+30) AS decimal(9,2)) / 3) AS char(1)), epa.epi_ReferralFaId
""")

con=connections("hchb")
df2=pd.read_sql(hchb_q2, con=con)
con.close()


df2['epi_ReferralFaId'] = df2['epi_ReferralFaId'].astype(str)
df2=df2.merge(referral_xwalk, left_on='epi_ReferralFaId', right_on='ID', how='left')
df2=df2.groupby(['Quarter', 'PARENT'])['denominator'].sum().reset_index()


df_h = df1.merge(df2, left_on=('Quarter', 'PARENT'), right_on=('Quarter','PARENT'), how='left')
df_h = df_h.sort_values(by=['Quarter']).reset_index(drop=True)
df_h = df_h[['Quarter', 'PARENT','numerator', 'denominator']]



# pcrs  rehosp(numerator):

pcrs_q1=(
"""
SELECT cf.case_num, cf.admission_date, cf.referral_ID, h1.hosp_dt, h1.hosp_dt- cf.admission_date AS hosp_days, CASE WHEN   h1.hosp_dt - cf.admission_date <= 30 THEN 1
ELSE 0 END AS hosp_30
from ms_owner.vw_case_facts cf
LEFT JOIN (  select cf1.case_num, cf1.admission_date, MIN(DECODE(ce.acute_hosp, 1, ce.end_dt, NULL)) AS hosp_dt
from ms_owner.vw_case_facts cf1
INNER JOIN oasis_owner.assessment am ON (am.case_num = cf1.case_num)
LEFT OUTER JOIN oasis_owner.cms_episode ce ON (cf1.case_num = ce.case_num)
WHERE ce.obqi_flag = 1
AND am.cms_epi_no = 1
AND am.cms_epi_seq_no = 1
AND cf1.admission_date >= """+"'"+o_start_30+"'""""
AND cf1.admission_date < """+"'"+o_stop_30+"'""""
AND cf1.REFERRAL_ID in """+str(referral_id)+"""
group by cf1.case_num, cf1.admission_date) h1
on h1.case_num = cf.case_num
where cf.admission_date >="""+"'"+o_start_30+"'""""
and cf.admission_date < """+"'"+o_stop_30+"'""""
and cf.referral_ID in """+str(referral_id)+"""
""")


or_con=connections("pcrs")
df1=pd.read_sql(pcrs_q1, con=or_con)
or_con.close()


df1 = df1[df1.HOSP_30 != 0].reset_index()
df1['adjusted_date'] = df1['ADMISSION_DATE'] + dt.timedelta(days = 30)
df1['Quarter']= df1['adjusted_date'].dt.year.astype(str) + '-' + 'Q' + df1['adjusted_date'].dt.quarter.astype(str)
df1=df1.merge(referral_xwalk, left_on='REFERRAL_ID', right_on='ID', how='left')
df1=df1.groupby(['Quarter', 'PARENT']).size().reset_index(name = 'numerator')


#pcrs admissions(denominator)

pcrs_q2="""
Select c.case_num, C.admission_date, c.referral_ID
from MS_OWNER.VW_CASE_FACTS c
where c.admission_date >= """+"'"+o_start_30+"'""""
and c.admission_date <  """+"'"+o_stop_30+"'""""
and c.referral_ID in"""+str(referral_id)+"""
and C.COMPANY_CODE='VNS'
and C.REFERRAL_TYPE='A'
"""
or_con=connections("pcrs")
df2=pd.read_sql(pcrs_q2, con=or_con)
or_con.close()


df2['adjusted_date'] = df2['ADMISSION_DATE'] + dt.timedelta(days = 30)
df2['Quarter']= df2['adjusted_date'].dt.year.astype(str) + '-' + 'Q' + df2['adjusted_date'].dt.quarter.astype(str)
df2=df2.merge(referral_xwalk, left_on='REFERRAL_ID', right_on='ID', how='left')
df2=df2.groupby(['Quarter', 'PARENT']).size().reset_index(name = 'denominator')


df_p = df1.merge(df2, left_on=('Quarter','PARENT'), right_on=('Quarter', 'PARENT'), how='left')
df_p = df_p.sort_values(by=['Quarter']).reset_index(drop=True)
df_p= df_p[['Quarter', 'PARENT','numerator', 'denominator']]

df_rehosp = pd.concat([df_h,df_p]).reset_index(drop=True)




#functions to create the template slide and displays the title below


def summaryX_setup(x):
    title_placeholder=x.shapes.title
    title_placeholder.text='CHHA Rehospitalization Trends:  '+str(system)+' System'
    return

def summaryY_setup(x):
    title_placeholder=x.shapes.title
    title_placeholder.text = 'CHHA Patient Admission distribution region:  '+str(system)+' System'
    return

def templateX_setup(x, i, cr):
    title_placeholder = x.shapes.title
    title_placeholder.text = 'CHHA Rehospitalization Trends:  '+cr
    return

def templateY_setup(x, i, cr):
    title_placeholder = x.shapes.title
    title_placeholder.text = 'CHHA Patient Admission distribution region:  '+cr
    return


#functions to create charts

def summaryX_chart(x):

    df=df_rehosp.groupby(['Quarter'])['numerator', 'denominator'].sum().reset_index()
    df['rate']= df['numerator']/df['denominator']
    m_date=df['Quarter'].astype(str).tolist()
    m_data=tuple(df['rate'])

    chart_data=ChartData()
    chart_data.categories=m_date
    chart_data.add_series('Rehospitalization',m_data, '0.0%')
    chart_title=str(b_reporting)+' - '+str(e_reporting)+' Rehospitalization Trend'
    chart_layout_line_p(x, chart_data, chart_title)
    return


def summaryY_chart(x):

    df=df_adm.groupby(['quarter', 'region']).size().reset_index(name = 'admissions by region')
    df = df[df['quarter'] == df['quarter'].max()].reset_index(drop=True)
    df['Admission percentage'] = df['admissions by region']/df['admissions by region'].sum()
    m_region=df['region'].astype(str).tolist()
    m_data=tuple(df['Admission percentage'])


    chart_data=ChartData()
    chart_data.categories=m_region
    chart_data.add_series('Admission Percentage Distribution', m_data, '0.0%')
    chart_title=str(b_reporting)+' -  '+str(e_reporting)+' Patient Region Distribution Breakout'
    chart_layout_pie(x,chart_data, chart_title)
    return



def templateX_chart(x,i):
    df = df_rehosp.groupby(['PARENT', 'Quarter'])['numerator', 'denominator'].sum().reset_index()
    df['rate']= df['numerator']/df['denominator']
    m_date=df.loc[df['PARENT']==i, 'Quarter'].values.tolist()
    m_rehospitalization=tuple(df.loc[df['PARENT']==i, 'rate'])

    chart_data=ChartData()
    chart_data.categories=m_date
    chart_data.add_series('Rehospitalization', m_rehospitalization, '0.0%')
    chart_title='Rehospitalization Trends '+str(b_reporting)+' - '+str(e_reporting)
    chart_layout_line_p(x,chart_data, chart_title)
    return



def templateY_chart(x,i):

   df=df_adm.groupby(['quarter', 'PARENT', 'region']).size().reset_index(name = 'admissions by region')
   df = df[df['quarter'] == df['quarter'].max()].reset_index(drop=True)
   dt = df.groupby(['PARENT'])['admissions by region'].sum().reset_index(name = 'Total admissions in site')
   df=df.merge(dt, left_on = 'PARENT', right_on = 'PARENT', how='inner')
   df['admission percentage'] = df['admissions by region']/df['Total admissions in site']

   m_region=df.loc[df['PARENT']==i, 'region'].values.tolist()
   m_data=tuple(df.loc[df['PARENT']==i, 'admission percentage'])

   chart_data=ChartData()
   chart_data.categories=m_region
   chart_data.add_series('admission Percentage Distribution', m_data, '0.0%')
   chart_title='Patient Region Distribution Breakout'+str(b_reporting)+' - '+str(e_reporting)
   chart_layout_pie(x,chart_data, chart_title)
   return



prs=Presentation()
title_slide=prs.slides.add_slide(prs.slide_layouts[2])
title=title_slide.shapes.title
title.text=str(system)



summaryX=prs.slides.add_slide(prs.slide_layouts[5])
summaryX_setup(summaryX)
summaryX_chart(summaryX)
summaryY=prs.slides.add_slide(prs.slide_layouts[5])
summaryY_setup(summaryY)
summaryY_chart(summaryY)


for row in df_locations.itertuples():
    i=row[1]
    cr=row[2]
    s=row[3]

    templateX=prs.slides.add_slide(prs.slide_layouts[5])
    templateX_setup(templateX, i, cr)
    comments(templateX)
    templateX_chart(templateX, i)
    templateY=prs.slides.add_slide(prs.slide_layouts[5])
    templateY_setup(templateY, i, cr)
    comments(templateY)
    templateY_chart(templateY,i)


prs.save('customer_reporting.pptx')
